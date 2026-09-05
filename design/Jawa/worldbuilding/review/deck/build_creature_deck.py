#!/usr/bin/env python3
"""Render the creature art register as a SPATIAL, drag-and-group PowerPoint deck.

Owner, 2026-09-05: *"An endless scroll list isn't human friendly... Are there
ways to more visually group and move these creatures perhaps using PowerPoint
as we did the tech tree?"*  This is the creature-art answer to the same loop
that worked for the research trees (build_tree_pptx.py -> he shuffles ->
owner_deck_answers_*.json read back).  It sits ALONGSIDE creature_register.html,
it does not replace it.

Reading the deck
----------------
* **Slide = one biome cluster** (the register's `group`).  A cluster too big for
  one slide is paginated — every page carries all four bands, so a drag never
  needs a slide change.
* **Four horizontal BANDS down the slide** = the decision:
  KEEP ART / REGENERATE / REGEN + RESCALE / CUT.  A creature starts in the band
  matching its current prefill (creature_register.decisions.json).
  **Dragging a creature into another band IS the decision.**  Deleting it = cut.
* Within a band, left -> right is smallest -> largest by true in-game draw size,
  so the size ladder stays visible while he works.
* Each creature is a GROUP shape (thumbnail + name caption) named
  ``CR:<defName>`` — one object, so the caption follows the drag and one Delete
  removes both.

Read back with ``read_creature_deck.py`` (writes the same row shape the research
deck used).  ``creature_deck_manifest.json`` beside the .pptx is what tells the
reader which creatures were SUPPOSED to be there, so a deletion is a fact and
not a silent drop.

Run from anywhere:
    python3 design/Jawa/worldbuilding/review/deck/build_creature_deck.py            # 3 biggest clusters
    python3 .../build_creature_deck.py --all                                        # every cluster
    python3 .../build_creature_deck.py --clusters 6 --thumb-px 128 --out foo.pptx
"""
import argparse
import hashlib
import json
import math
import os
import sys
from collections import Counter, OrderedDict

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
REVIEW = os.path.dirname(HERE)
ROWS = os.path.join(REVIEW, "creature_register_rows.json")
DECISIONS = os.path.join(REVIEW, "creature_register.decisions.json")
ART = os.path.join(REVIEW, "creature_art")
THUMBS = os.path.join(HERE, "thumbs")
OUT = os.path.join(HERE, "creature_deck.pptx")
MANIFEST = os.path.join(HERE, "creature_deck_manifest.json")

# ------------------------------------------------------------------ bands
# key, title, fill (band zone tint), accent (strip), what it means
BANDS = [
    ("keep",    "KEEP ART",        "e7f1e3", "4f8a55",
     "the shipping art is good enough — no work queued"),
    ("regen",   "REGENERATE",      "fdf0d8", "c08a2e",
     "art must be redrawn; the size on the map is right"),
    ("rescale", "REGEN + RESCALE", "fbe6dd", "c2632e",
     "redraw AND change drawSize — it renders at the wrong scale"),
    ("cut",     "CUT",             "efe3e6", "9c3f52",
     "the creature leaves the planet entirely (Cherry Picker)"),
]
BAND_KEYS = [b[0] for b in BANDS]
BAND_TITLE = {b[0]: b[1] for b in BANDS}
BAND_FILL = {b[0]: b[2] for b in BANDS}
BAND_ACCENT = {b[0]: b[3] for b in BANDS}
BAND_MEANS = {b[0]: b[4] for b in BANDS}

# ------------------------------------------------------------------ geometry
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.08)
HEADER_H = Inches(0.44)
STRIP_H = Inches(0.20)          # band label strip, inside the band zone
BAND_GAP = Inches(0.04)
THUMB = Inches(0.60)            # picture edge
CAP_H = Inches(0.17)            # caption under it
CELL_W = Inches(0.68)
CELL_H = THUMB + CAP_H
ROW_PITCH = CELL_H + Inches(0.05)
COLS = int((SLIDE_W - 2 * MARGIN) // CELL_W)
# rows of thumbnails available on one slide, shared out between the four bands
BODY_H = SLIDE_H - HEADER_H - 2 * MARGIN - 4 * STRIP_H - 3 * BAND_GAP
ROWS_PER_SLIDE = int(BODY_H // ROW_PITCH)


def rgb(h):
    return RGBColor.from_string(h)


def no_autofit(tf):
    el = tf._txBody
    bp = el.find(qn("a:bodyPr"))
    for tag in ("a:normAutofit", "a:spAutoFit"):
        e = bp.find(qn(tag))
        if e is not None:
            bp.remove(e)
    bp.append(el.makeelement(qn("a:noAutofit"), {}))


def text_in(shape_or_box, text, pt, bold=False, color="2b2418",
            align=PP_ALIGN.CENTER):
    tf = shape_or_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(1)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(pt)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    no_autofit(tf)
    return shape_or_box


def rect(shapes, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    s = shapes.add_shape(shape, int(x), int(y), int(w), int(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    if line:
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


# ------------------------------------------------------------------ thumbnails
def thumb_for(defname, art, px):
    """Downscale the DETAIL render (the art is what is being judged)."""
    src = art.get("detail")
    if not src:
        return None
    src = os.path.join(REVIEW, src) if not os.path.isabs(src) else src
    if not os.path.exists(src):
        return None
    os.makedirs(THUMBS, exist_ok=True)
    dst = os.path.join(THUMBS, "%s.%d.jpg" % (defname, px))
    if not os.path.exists(dst):
        im = Image.open(src).convert("RGB")
        im.thumbnail((px, px), Image.LANCZOS)
        im.save(dst, "JPEG", quality=82, optimize=True)
    return dst


# ------------------------------------------------------------------ data
def true_size(row):
    ds = row.get("drawSize") or [1, 1]
    try:
        d = max(float(ds[0]), float(ds[1]))
    except (TypeError, ValueError, IndexError):
        d = 1.0
    try:
        b = float(row.get("bodySize") or 0)
    except (TypeError, ValueError):
        b = 0.0
    return (d, b)


def load():
    reg = json.load(open(ROWS, encoding="utf-8"))
    dec = json.load(open(DECISIONS, encoding="utf-8"))["decisions"]
    rows = []
    for r in reg["rows"]:
        d = dec.get(r["defName"]) or {}
        # The owner's own decision wins over the generated prefill.
        # ⚠️ r["cut"] is NOT a review decision — it is Cherry Picker's CURRENT
        # state (270 creatures are already cut from the live game).  It is a
        # marker on the tile, never a band placement, or the deck would read
        # back 270 fresh "cut" rulings the owner never made.
        band = d.get("decision") or d.get("prefill") or "keep"
        if band not in BAND_KEYS:
            band = "keep"
        rows.append({
            "defName": r["defName"],
            "label": r.get("label") or r["defName"],
            "biome": r.get("group") or "(ungrouped)",
            "band": band,
            "size": true_size(r),
            "alreadyCut": bool(r.get("cut")),
            "art": r.get("art") or {},
        })
    return reg["meta"], rows


def paginate(members):
    """Split one cluster into pages; every page carries all four bands.

    Returns a list of pages, each ``{band_key: [rows in x-order]}``.
    """
    left = OrderedDict((k, list(members.get(k, []))) for k in BAND_KEYS)
    pages = []
    while any(left.values()) or not pages:
        # every band gets at least one row so it is always a drop target
        alloc = {k: 1 for k in BAND_KEYS}
        spare = ROWS_PER_SLIDE - len(BAND_KEYS)
        need = {k: max(0, math.ceil(len(v) / COLS) - 1) for k, v in left.items()}
        while spare > 0 and any(need.values()):
            k = max(need, key=lambda k: (need[k], len(left[k])))
            alloc[k] += 1
            need[k] -= 1
            spare -= 1
        page = {}
        for k in BAND_KEYS:
            take = alloc[k] * COLS
            page[k] = left[k][:take]
            left[k] = left[k][take:]
        pages.append((page, alloc))
        if not any(left.values()):
            break
    return pages


# ------------------------------------------------------------------ slides
def legend_slide(prs, meta, clusters, total, scope_note):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide.shapes, 0, 0, SLIDE_W, Inches(0.9), "2b2418")
    t = slide.shapes.add_textbox(Inches(0.3), Inches(0.02), SLIDE_W - Inches(0.6), Inches(0.86))
    tf = t.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Creature art review — drag a creature into the band you want"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = rgb("e8dcc0")
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = ("One slide per biome cluster.  %s   ·   register built %s"
               % (scope_note, meta.get("builtUtc", "?")))
    r2.font.size = Pt(11)
    r2.font.color.rgb = rgb("b3a68d")
    no_autofit(tf)

    y = Inches(1.15)
    h = Inches(0.72)
    for k, title, fill, accent, means in BANDS:
        rect(slide.shapes, Inches(0.35), y, SLIDE_W - Inches(0.7), h, fill, accent)
        lab = slide.shapes.add_textbox(Inches(0.5), y, Inches(2.6), h)
        text_in(lab, title, 16, bold=True, color=accent, align=PP_ALIGN.LEFT)
        txt = slide.shapes.add_textbox(Inches(3.2), y, SLIDE_W - Inches(3.9), h)
        text_in(txt, means, 12, color="4a4238", align=PP_ALIGN.LEFT)
        y += h + Inches(0.10)

    note = slide.shapes.add_textbox(Inches(0.35), y + Inches(0.05),
                                    SLIDE_W - Inches(0.7), Inches(1.5))
    tf = note.text_frame
    tf.word_wrap = True
    lines = [
        "HOW TO USE IT — drag a creature's tile from the band it is in into the band you want. "
        "That drag IS the decision; nothing else needs typing.",
        "DELETING a tile (select it, press Delete) also means CUT — same as dropping it in the CUT band.",
        "A RED-OUTLINED tile with a \u2702 on its name is ALREADY cut from the live game by Cherry Picker. "
        "That is the current state, not a ruling — leave it where it sits to keep it, or move it like any other tile.",
        "Each tile is one grouped object: the picture and its name move together. Within a band, "
        "left → right runs smallest → largest at true in-game draw size, so the size ladder stays readable.",
        "A big cluster runs over several slides — every page carries all four bands, so a drag never needs a slide change.",
        "Tiles start in the band matching the CURRENT decision, so an untouched deck reads back unchanged. "
        "%d creatures across %d cluster(s) on this deck." % (total, clusters),
        "Read it back with:  python3 design/Jawa/worldbuilding/review/deck/read_creature_deck.py <your.pptx>",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = "•  " + line
        r.font.size = Pt(11)
        r.font.color.rgb = rgb("3d372e" if i < 2 else "5a5147")
        r.font.bold = i < 2
    no_autofit(tf)
    return slide


def creature_tile(shapes, x, y, row, thumb_px):
    """A group: thumbnail + caption, named CR:<defName>."""
    g = shapes.add_group_shape()
    g.name = "CR:" + row["defName"]
    img = thumb_for(row["defName"], row["art"], thumb_px)
    if img:
        g.shapes.add_picture(img, int(x + (CELL_W - THUMB) / 2), int(y),
                             int(THUMB), int(THUMB))
    else:
        ph = rect(g.shapes, x + (CELL_W - THUMB) / 2, y, THUMB, THUMB,
                  "d9d2c6", "8a8073")
        text_in(ph, "no art", 7, color="6b6155")
    if row.get("alreadyCut"):
        # already cut from the live game by Cherry Picker — a marker, not a ruling
        mark = rect(g.shapes, x + (CELL_W - THUMB) / 2, y, THUMB, THUMB,
                    "ffffff", "9c3f52")
        mark.fill.background()
        mark.line.width = Pt(1.5)
    cap = g.shapes.add_textbox(int(x), int(y + THUMB), int(CELL_W), int(CAP_H))
    lab = row["label"]
    pt = 6.0 if len(lab) <= 16 else (5.5 if len(lab) <= 24 else 5.0)
    if row.get("alreadyCut"):
        text_in(cap, "\u2702 " + lab, pt, color="9c3f52")
    else:
        text_in(cap, lab, pt, color="2b2418")
    return g


def cluster_slides(prs, cluster, members, thumb_px, manifest):
    pages = paginate(members)
    n_total = sum(len(v) for v in members.values())
    for pi, (page, alloc) in enumerate(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        rect(slide.shapes, 0, 0, SLIDE_W, HEADER_H, "2b2418")
        hdr = slide.shapes.add_textbox(Inches(0.15), 0, SLIDE_W - Inches(0.3), HEADER_H)
        # 🔑 the reader finds the slide's cluster from THIS shape's name.
        hdr.name = "SLIDEMETA:%s|%d/%d" % (cluster, pi + 1, len(pages))
        tf = hdr.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = cluster + ("   (%d/%d)" % (pi + 1, len(pages)) if len(pages) > 1 else "")
        r.font.size = Pt(17)
        r.font.bold = True
        r.font.color.rgb = rgb("e8dcc0")
        r2 = p.add_run()
        r2.text = "     %d creatures on this page of %d in the cluster    ·    drag between bands to decide" % (
            sum(len(v) for v in page.values()), n_total)
        r2.font.size = Pt(10)
        r2.font.color.rgb = rgb("a89a80")
        no_autofit(tf)

        y = HEADER_H + MARGIN
        for k in BAND_KEYS:
            rows_here = alloc[k]
            zone_h = STRIP_H + rows_here * ROW_PITCH
            zone = rect(slide.shapes, MARGIN, y, SLIDE_W - 2 * MARGIN, zone_h,
                        BAND_FILL[k], BAND_ACCENT[k])
            # 🔑 the reader gets the band boundaries from THIS shape's geometry.
            zone.name = "BAND:" + k
            strip = rect(slide.shapes, MARGIN, y, SLIDE_W - 2 * MARGIN, STRIP_H,
                         BAND_ACCENT[k])
            lab = slide.shapes.add_textbox(MARGIN + Inches(0.06), y,
                                           SLIDE_W - 2 * MARGIN, STRIP_H)
            text_in(lab, "%s   —   %s   (%d here)"
                    % (BAND_TITLE[k], BAND_MEANS[k], len(page[k])),
                    9, bold=True, color="ffffff", align=PP_ALIGN.LEFT)
            for i, row in enumerate(page[k]):
                col, rw = i % COLS, i // COLS
                cx = MARGIN + col * CELL_W
                cy = y + STRIP_H + rw * ROW_PITCH
                creature_tile(slide.shapes, cx, cy, row, thumb_px)
                manifest["creatures"][row["defName"]] = {
                    "label": row["label"],
                    "biome": cluster,
                    "slide": len(prs.slides._sldIdLst),   # 1-based index
                    "band": k,
                    "alreadyCut": row.get("alreadyCut", False),
                }
            y += zone_h + BAND_GAP
    return len(pages)


# ------------------------------------------------------------------ main
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--clusters", type=int, default=3,
                    help="how many of the largest biome clusters to render (default 3)")
    ap.add_argument("--all", action="store_true", help="render every cluster")
    ap.add_argument("--thumb-px", type=int, default=112)
    ap.add_argument("--out", default=OUT)
    ap.add_argument("--manifest", default=None)
    a = ap.parse_args()

    meta, rows = load()
    by_cluster = OrderedDict()
    for r in rows:
        by_cluster.setdefault(r["biome"], []).append(r)
    order = sorted(by_cluster, key=lambda c: (-len(by_cluster[c]), c))
    chosen = order if a.all else order[:a.clusters]

    total = sum(len(by_cluster[c]) for c in chosen)
    scope = ("all %d clusters" % len(order) if a.all
             else "PROTOTYPE: the %d largest of %d clusters" % (len(chosen), len(order)))

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    manifest = {
        "meta": {
            "generator": "build_creature_deck.py 1.0",
            "registerBuiltUtc": meta.get("builtUtc"),
            "deck": os.path.basename(a.out),
            "scope": scope,
            "thumbPx": a.thumb_px,
            "bands": BAND_KEYS,
            "what": ("Which creature was placed on which slide in which band when the "
                     "deck was generated. read_creature_deck.py diffs the owner's "
                     "edited deck against this, so a DELETED creature is a fact and "
                     "never a silent drop."),
        },
        "creatures": {},
    }
    legend_slide(prs, meta, len(chosen), total, scope)
    for c in chosen:
        members = {}
        for r in sorted(by_cluster[c], key=lambda r: (r["size"], r["label"])):
            members.setdefault(r["band"], []).append(r)
        cluster_slides(prs, c, members, a.thumb_px, manifest)

    prs.save(a.out)
    json.dump(manifest, open(a.manifest or MANIFEST, "w", encoding="utf-8"),
              indent=1, ensure_ascii=False)

    placed = len(manifest["creatures"])
    if placed != total:
        print("MISMATCH: %d of %d creatures did not land on a slide" % (total - placed, total))
        sys.exit(1)
    size = os.path.getsize(a.out)
    print("wrote %s" % a.out)
    print("  %s: %d creatures, %d slides (1 legend + %d cluster pages), %.2f MB"
          % (scope, placed, len(prs.slides._sldIdLst),
             len(prs.slides._sldIdLst) - 1, size / 1e6))
    print("  grid: %d cols x %d thumbnail rows per slide, %dpx thumbs" %
          (COLS, ROWS_PER_SLIDE, a.thumb_px))
    per_slide = []
    for s in prs.slides:
        per_slide.append(len(s.shapes))
    print("  shapes/slide: min %d, max %d, mean %.1f (legend %d)"
          % (min(per_slide), max(per_slide), sum(per_slide) / len(per_slide), per_slide[0]))
    print("  bands: " + ", ".join("%s %d" % (k, sum(1 for v in manifest["creatures"].values()
                                                    if v["band"] == k)) for k in BAND_KEYS))


if __name__ == "__main__":
    main()
