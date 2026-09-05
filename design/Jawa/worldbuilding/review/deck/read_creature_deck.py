#!/usr/bin/env python3
"""Read the owner's edited creature deck back into decision rows.

A view that cannot be captured is useless, so this is the other half of
``build_creature_deck.py``.  It opens the .pptx the owner shuffled and, per
creature, resolves:

  * which slide it now sits on            -> ``biome`` (the cluster)
  * which BAND its CENTRE now falls in    -> ``decision`` (keep/regen/rescale/cut)
  * whether the tile is gone entirely     -> ``removed`` (which also means cut)
  * where it sits left-to-right in that band -> ``x_order``

Row shape matches the research-deck precedent (owner_deck_answers_20260904.json).

🔴 It never silently drops a creature.  Every defName in the manifest lands in
``rows`` exactly once, or is named in ``meta.unresolved`` with the reason.  The
exit code is non-zero if anything is unresolved.

    python3 read_creature_deck.py <edited.pptx> [--manifest ...] [--out ...]
"""
import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
MANIFEST = os.path.join(HERE, "creature_deck_manifest.json")
OUT = os.path.join(HERE, "creature_deck_answers.json")

# A centre landing in the gap between two bands (or a hair outside the top or
# bottom one) snaps to the nearest band and is FLAGGED, never guessed silently.
SNAP_TOL = Emu(int(Inches(0.30)))


def collect(pptx_path):
    """-> (slides, tiles) with every CR: shape found anywhere in the deck."""
    prs = Presentation(pptx_path)
    slides, tiles = [], []
    for idx, slide in enumerate(prs.slides, start=1):
        cluster, page = None, None
        bands = []
        for sh in slide.shapes:
            name = sh.name or ""
            if name.startswith("SLIDEMETA:"):
                rest = name[len("SLIDEMETA:"):]
                cluster, _, page = rest.partition("|")
            elif name.startswith("BAND:"):
                bands.append({"key": name[len("BAND:"):],
                              "top": int(sh.top), "bottom": int(sh.top + sh.height)})
            elif name.startswith("CR:"):
                tiles.append({
                    "defName": name[len("CR:"):],
                    "slide": idx,
                    "cx": int(sh.left + sh.width / 2),
                    "cy": int(sh.top + sh.height / 2),
                })
        bands.sort(key=lambda b: b["top"])
        slides.append({"index": idx, "cluster": cluster, "page": page, "bands": bands})
    return slides, tiles


def band_for(bands, cy):
    """-> (band_key, snapped) or (None, reason)."""
    for b in bands:
        if b["top"] <= cy < b["bottom"]:
            return b["key"], False
    if not bands:
        return None, "slide has no bands"
    best, dist = None, None
    for b in bands:
        d = 0 if b["top"] <= cy < b["bottom"] else min(abs(cy - b["top"]),
                                                       abs(cy - b["bottom"]))
        if dist is None or d < dist:
            best, dist = b, d
    if dist <= SNAP_TOL:
        return best["key"], True
    return None, "centre %.2fin from the nearest band" % (dist / 914400.0)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--manifest", default=MANIFEST)
    ap.add_argument("--out", default=OUT)
    ap.add_argument("--quiet", action="store_true")
    a = ap.parse_args()

    man = json.load(open(a.manifest, encoding="utf-8"))
    expected = man["creatures"]
    slides, tiles = collect(a.pptx)
    by_slide = {s["index"]: s for s in slides}

    unresolved, duplicates, unknown, snapped = [], [], [], []
    seen = {}
    for t in tiles:
        dn = t["defName"]
        s = by_slide[t["slide"]]
        if dn in seen:
            duplicates.append({"defName": dn, "slide": t["slide"],
                               "kept": seen[dn]["slide"]})
            continue
        if dn not in expected:
            unknown.append({"defName": dn, "slide": t["slide"]})
        key, flag = band_for(s["bands"], t["cy"])
        if key is None:
            unresolved.append({"defName": dn, "slide": t["slide"],
                               "cluster": s["cluster"], "reason": flag})
            continue
        if flag:
            snapped.append(dn)
        seen[dn] = {"slide": t["slide"], "cluster": s["cluster"], "band": key,
                    "cx": t["cx"], "snapped": bool(flag)}

    # x_order is the rank of the centre within its own (slide, band)
    order_keys = {}
    for dn, v in sorted(seen.items(), key=lambda kv: (kv[1]["slide"], kv[1]["band"],
                                                      kv[1]["cx"], kv[0])):
        k = (v["slide"], v["band"])
        order_keys[k] = order_keys.get(k, -1) + 1
        v["x_order"] = order_keys[k]

    rows, changed, removed_names = [], [], []
    for dn, meta in expected.items():
        got = seen.get(dn)
        if got is None:
            if any(u["defName"] == dn for u in unresolved):
                continue                      # already reported, not a deletion
            rows.append({"defName": dn, "label": meta["label"],
                         "biome": meta["biome"], "decision": "cut",
                         "removed": True, "x_order": None,
                         "was": meta["band"], "note": "tile deleted from the deck"})
            removed_names.append(dn)
            if meta["band"] != "cut":
                changed.append({"defName": dn, "from": meta["band"], "to": "cut (deleted)"})
            continue
        row = {"defName": dn, "label": meta["label"],
               "biome": got["cluster"] or meta["biome"],
               "decision": got["band"], "removed": False,
               "x_order": got["x_order"], "was": meta["band"]}
        if got["snapped"]:
            row["snapped"] = True
        if got["cluster"] and got["cluster"] != meta["biome"]:
            row["note"] = "moved to the %s slide" % got["cluster"]
        rows.append(row)
        if got["band"] != meta["band"]:
            changed.append({"defName": dn, "label": meta["label"],
                            "from": meta["band"], "to": got["band"]})

    out = {
        "meta": {
            "reader": "read_creature_deck.py 1.0",
            "deck": os.path.abspath(a.pptx),
            "manifest": os.path.abspath(a.manifest),
            "authority": ("OWNER — the band each tile now sits in is his ruling. "
                          "removed=true is a deleted tile and means cut."),
            "expected": len(expected),
            "resolved": len(seen),
            "rows": len(rows),
            "deleted": len(removed_names),
            "changed": len(changed),
            "snapped": snapped,
            "unresolved": unresolved,
            "duplicates": duplicates,
            "unknown": unknown,
        },
        "changes": changed,
        "rows": rows,
    }
    json.dump(out, open(a.out, "w", encoding="utf-8"), indent=1, ensure_ascii=False)

    if not a.quiet:
        print("read %s" % a.pptx)
        print("  %d expected, %d rows written, %d deleted, %d unresolved, "
              "%d duplicate, %d unknown, %d snapped"
              % (len(expected), len(rows), len(removed_names), len(unresolved),
                 len(duplicates), len(unknown), len(snapped)))
        print("  %d decision(s) CHANGED from the generated deck:" % len(changed))
        for c in changed[:40]:
            print("    %-34s %s -> %s" % (c["defName"], c["from"], c["to"]))
        if len(changed) > 40:
            print("    ... and %d more" % (len(changed) - 40))
        print("  wrote %s" % a.out)

    accounted = len(rows) + len(unresolved)
    if accounted != len(expected):
        print("ACCOUNTING FAILURE: %d expected but %d accounted for"
              % (len(expected), accounted))
        return 2
    if unresolved:
        print("UNRESOLVED: %d tile(s) could not be placed in a band — see meta.unresolved"
              % len(unresolved))
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
