#!/usr/bin/env python3
"""Round-trip proof: drag two creatures between bands the way the owner would,
save a copy, and let read_creature_deck.py find exactly those two changes.

This is the deliverable's real test — a deck whose edits cannot be read back is
just a picture.  It is a regression harness, not part of the review loop.

    python3 simulate_owner_edit.py [--deck creature_deck.pptx]
"""
import argparse
import os
import sys

from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--deck", default=os.path.join(HERE, "creature_deck.pptx"))
    ap.add_argument("--out", default=os.path.join(HERE, "creature_deck_simulated_edit.pptx"))
    a = ap.parse_args()

    prs = Presentation(a.deck)
    slide = list(prs.slides)[1]              # first cluster slide
    bands = {}
    tiles = []
    for sh in slide.shapes:
        n = sh.name or ""
        if n.startswith("BAND:"):
            bands[n[5:]] = sh
        elif n.startswith("CR:"):
            tiles.append(sh)

    # two tiles that currently sit in KEEP, dropped into REGEN and CUT
    keep = bands["keep"]
    in_keep = [t for t in tiles
               if keep.top <= t.top + t.height / 2 < keep.top + keep.height]
    moves = [(in_keep[0], "regen"), (in_keep[-1], "cut")]

    done = []
    for tile, target in moves:
        b = bands[target]
        tile.left = int(b.left + Inches(0.5))
        tile.top = int(b.top + Inches(0.30))   # below the band's label strip
        done.append((tile.name[3:], target))

    prs.save(a.out)
    print("simulated the owner's drags on slide 2 of %s" % os.path.basename(a.deck))
    for dn, target in done:
        print("   %-34s -> %s" % (dn, target))
    print("wrote %s" % a.out)
    return 0


if __name__ == "__main__":
    sys.exit(main())
