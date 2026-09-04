#!/usr/bin/env python3
"""Render the v4 tree proposal as a SHUFFLEABLE PowerPoint deck.

Owner's request, 2026-09-04: *"a Powerpoint presentation with each technology
as a small box I can shuffle around. Each of the current tech trees should be
their own slide. Array the little boxes in the order of technology bin
development left (primitive) to right (advanced)."* Revised same day:
*"all the technologies an identically defined width (smallest of current
options). Recolor each box according to the kind of thing it is."*

Round 3 (2026-09-04): layers owner_deck_answers_20260904.json (the owner's
shuffled deck, extracted — HIS trees/tiers/removals, authoritative) and
deck_round2_adjustments.json (BENCH tier spreads on the bunched trees) over
restructured_model_v4.json, then writes research_trees_boxes.pptx.
Border language: THICK border = an owner placement (his move), DASHED = a
BENCH round-2 tier spread, thin solid = unmoved from v4. Each research row is one INDIVIDUAL shape —
movable, editable, deletable in PowerPoint or Google Slides. Boxes are all one
width, sit in five tier columns (T0 primitive … T4 advanced, cost-sorted
within a tier), and are colored by CATEGORY of what they unlock (weapon,
apparel, vehicles, power, droids, …) — slide 1 is the legend. Category comes
from keyword rules over label+unlocks with a per-tree fallback; it is a
review aid, not a ruling — the owner corrects by recoloring a box.

⚠️ Once the owner has shuffled the deck, the .pptx holds HIS decisions and this
generator must not be re-run over it — write to a new filename instead.

Run from repo root:  python3 design/Jawa/research_review/build_tree_pptx.py
"""
import json
import os
import re
import sys
from collections import Counter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE_DASH
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
MODEL = os.path.join(HERE, "restructured_model_v4.json")
ANSWERS = os.path.join(HERE, "owner_deck_answers_20260904.json")
ADJUST = os.path.join(HERE, "deck_round2_adjustments.json")
OUT = os.path.join(HERE, "research_trees_boxes_round3.pptx")

# Slide order and accents — same roster as build_tree_visual_v4.py.
TREES = [
    ("Scavenger", "c2a06a", "the pride-free floor — fire, water, food, hide, door, trap"),
    ("The Hearth", "b98a4e", "comfort & culture — cooking, brew, furniture, art, music, cloth"),
    ("The Refinery", "a8764e", "what sand and wreck become — fuels, chems, drugs, ores, synthetics"),
    ("The Workshop", "bd6f4e", "making & mending — smithing, machining, electronics, vehicles, power"),
    ("Powder & Slug", "8a8a6a", "kills by MASS — guns, cannon, mortars, blades, the Watch"),
    ("Blasterworks", "c25a4a", "kills by HEAT — the blaster spine, plasma, beam, tibanna"),
    ("The Strange Schools", "9c6b8a", "kills by STRANGER physics — ion, sonic, vibro, relics, saber"),
    ("The Shell", "7089a0", "not dying — armors, shields, warcaskets, the maker doctrines"),
    ("Droidsmith", "6f9083", "the scavenger's honest skill — repair, reconstruction, maintenance"),
    ("The Waking Mind", "5f8a74", "minds you make and minds you bind — the AI ladder"),
    ("THE SHIP", "4a7d86", "the Utinni herself — gravtech, her systems, her guns, her memory"),
    ("The Reach", "8a6b9c", "the trap — flesh, genes, bionics, archotech, priced brutally"),
    # Locked (faction-gated) trees follow the twelve.
    ("The Junker Yards", "a8503f", "LOCKED · Junkers — everything warcasket; raid loot and quests only"),
    ("The Foundry Hive", "7d9c5a", "LOCKED · Foundry Hive — sonic, hivetech, battle droids; they trade"),
    ("The Ascendant Ladder", "8a6b9c", "LOCKED · Helix — genes, biosculpting, the flesh ladder"),
    ("The Unbolting", "6f9083", "LOCKED · Free Droid Enclaves — building droids at all"),
]

TIERS = ["T0", "T1", "T2", "T3", "T4"]
TIER_LABEL = {"T0": "T0 · ≤600", "T1": "T1 · 600–1600",
              "T2": "T2 · 1600–3000", "T3": "T3 · 3000–5000",
              "T4": "T4 · 5000+"}
# Bin gradient for the COLUMN HEADERS only — boxes are colored by category.
TIER_FILL = {"T0": "f2ead9", "T1": "e8dcc0", "T2": "ddcda6",
             "T3": "cdb989", "T4": "bda36c"}

# ---------------------------------------------------------------- categories
# (name, fill, keyword regex) — first match wins, tested against
# "label | unlock defNames" lowercased. Fallback per tree below.
CATEGORIES = [
    ("droids",       "7fcfc3", r"droid|mechanoid|positronic|automaton|robot|astromech|protocol|servitor|mech_"),
    ("space",        "92b4e3", r"gravship|grav |gravtech|grav_|orbital|starflight|vacuum|hyperdrive|starship|shuttle|spacedrive|_ship|ship_"),
    ("vehicles",     "f0a860", r"vehicle|blueprint|speeder|rover|wheel|aircraft|helicopter|boat|raft|caravan"),
    ("security",     "a8b8a0", r"turret|trap|ied|fortif|sandbag|barricade|dugout|trench|defen[cs]e|firefoam|foam pop"),
    ("weapon",       "e59896", r"weapon|gun|rifle|pistol|blaster|cannon|mortar|sword|blade|saber|sabre|bow\b|spear|melee|ammo|ammunition|grenade|launcher|vibro|slug|revolver|carbine|smg|shotgun|minigun|staff|baton|electrostaff|disruptor"),
    ("apparel",      "c7a3d6", r"armor|armour|apparel|clothing|tailor|vest|helmet|warcasket|shield|suit\b|garb|robe|fatigues|plate\b|mask\b"),
    ("power",        "f2dd7c", r"power|electric|battery|batteries|solar|geothermal|watermill|reactor|generator|conduit|turbine|dynamo"),
    ("biological",   "9ecf8e", r"bionic|prosthe|gene\b|genes|geneti|xeno|biosculpt|medicine|medical|hospital|surgery|implant|organ|drug|penoxy|neutroamine|psychite|wake-up|go-juice|serum|healroot|anesthe"),
    ("industry",     "c4a884", r"machining|fabricat|smith|smelt|refin|chemfuel|chemistry|processor|factory|assembly|milling|drill|mining|deep drill|extractor|processing|kiln|forge|foundry"),
    ("culture",      "e3a8c4", r"art\b|artistic|music|instrument|furniture|recreation|joy\b|sculpt|decorat|carpet|floor|drapes|lamp|brazier|television|chess|jewel|sauna|hot ?tub|curtain|torch|lighting"),
    ("life_support", "f0d9b8", r"nutrition|cooking|stove|meal|farming|agricultur|hydroponic|crop|plant_|plant |brewing|brewery|preserv|freezer|refrigerat|water|well\b|fire\b|firecraft|bedroll|heating|cooling|cooler|air condition|survival|kibble|dispenser|cultivat|animal|devilstrand|cocoa|grill|fry|cheese|bathroom|shower|sink\b|door|gate|bridge|storage|urn\b|paste|moisture"),
    ("equipment",    "d8c8b8", r"tool|gadget|pack\b|utility|comms|communicat|sensor|scanner|analyzer|radar|goggles|binocular|tinker|hack|vitals"),
]
GENERAL = ("general", "c8c8c8")
TREE_FALLBACK = {
    "Powder & Slug": "weapon", "Blasterworks": "weapon",
    "The Strange Schools": "weapon", "The Shell": "apparel",
    "Droidsmith": "droids", "The Unbolting": "droids",
    "The Waking Mind": "droids", "THE SHIP": "space",
    "The Refinery": "industry", "The Junker Yards": "apparel",
    "The Foundry Hive": "droids", "The Ascendant Ladder": "biological",
}
CAT_FILL = {n: f for n, f, *_ in CATEGORIES}
CAT_FILL[GENERAL[0]] = GENERAL[1]
_RULES = [(n, re.compile(rx)) for n, _f, rx in CATEGORIES]


def categorize(m):
    hay = (m["label"] + " | " + " ".join(m.get("unlocks") or [])).lower()
    for name, rx in _RULES:
        if rx.search(hay):
            return name
    return TREE_FALLBACK.get(m["tab4"], GENERAL[0])


# ---------------------------------------------------------------- geometry
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
HEADER_H = Inches(0.72)
TIERBAR_H = Inches(0.26)
MARGIN = Inches(0.12)
BOX_H = Emu(int(Inches(0.30)))
VGAP = Emu(int(Inches(0.045)))
HGAP = Emu(int(Inches(0.06)))
COL_W = (SLIDE_W - 2 * MARGIN) // 5
# 🔑 ONE width for every box — the smallest of v1's options (its 3-across case).
SUB_COLS = 3
BOX_W = (COL_W - HGAP - (SUB_COLS - 1) * HGAP) // SUB_COLS


def rgb(hexstr):
    return RGBColor.from_string(hexstr)


def tier_of(m):
    return m.get("_tier") or m.get("tier4") or m.get("tier3") or "T0"


def cost_of(m):
    for k in ("cost4", "cost2", "cost"):
        v = m.get(k)
        if v not in (None, ""):
            try:
                return float(v)
            except (TypeError, ValueError):
                pass
    return 0.0


def no_autofit(tf):
    # Keep the shape the size we set; PowerPoint's autofit would fight shuffling.
    el = tf._txBody
    bodyPr = el.find(qn("a:bodyPr"))
    for tag in ("a:normAutofit", "a:spAutoFit"):
        e = bodyPr.find(qn(tag))
        if e is not None:
            bodyPr.remove(e)
    bodyPr.append(el.makeelement(qn("a:noAutofit"), {}))


def add_box(slide, x, y, w, h, text, fill, line, font_pt, bold=False,
            font_color="2b2418", align=PP_ALIGN.CENTER, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(shape, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = rgb(fill)
    box.line.color.rgb = rgb(line)
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_pt)
    r.font.bold = bold
    r.font.color.rgb = rgb(font_color)
    no_autofit(tf)
    return box


def font_for(label):
    n = len(label)
    return 8 if n <= 14 else (7 if n <= 28 else 6)


def legend_slide(prs, counts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, SLIDE_W, HEADER_H, "", "2b2418", "2b2418", 10,
            shape=MSO_SHAPE.RECTANGLE)
    title = slide.shapes.add_textbox(MARGIN, Emu(0), SLIDE_W - 2 * MARGIN, HEADER_H)
    tf = title.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Legend — box color = what the tech unlocks"
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = rgb("d8cfc0")
    no_autofit(tf)
    names = [n for n, _f, _rx in CATEGORIES] + [GENERAL[0]]
    y0 = HEADER_H + Inches(0.35)
    chip_w, chip_h = Inches(1.5), Inches(0.4)
    for i, name in enumerate(names):
        col, row = divmod(i, 7)
        x = MARGIN + Inches(0.3) + col * Inches(4.2)
        y = y0 + row * (chip_h + Inches(0.22))
        add_box(slide, x, y, chip_w, chip_h, name, CAT_FILL[name], "6b5f4a", 12)
        lbl = slide.shapes.add_textbox(x + chip_w + Inches(0.15), y,
                                       Inches(2.4), chip_h)
        tf = lbl.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = "%d techs" % counts.get(name, 0)
        r.font.size = Pt(11); r.font.color.rgb = rgb("5a5147")
        no_autofit(tf)
    note = slide.shapes.add_textbox(MARGIN + Inches(0.3), SLIDE_H - Inches(0.9),
                                    SLIDE_W - Inches(1), Inches(0.6))
    p = note.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = ("Columns run T0 (primitive) left → T4 (advanced) right, cost-sorted "
              "within a tier. Colors are keyword-derived — recolor a box to correct it.")
    r.font.size = Pt(11); r.font.color.rgb = rgb("5a5147")
    no_autofit(note.text_frame)


def build():
    M = json.load(open(MODEL, encoding="utf-8"))
    ans = {r["defName"]: r for r in
           json.load(open(ANSWERS, encoding="utf-8"))["rows"]}
    adj = json.load(open(ADJUST, encoding="utf-8"))
    adj_by = {(a["tree"], a["label"]): a["to"] for a in adj["tier_moves"]}
    surv = []
    for m in M:
        if not m.get("tab4"):
            continue
        a = ans.get(m["defName"])
        if a is None:
            raise SystemExit("no owner answer for %s" % m["defName"])
        if a["removed"]:
            continue                      # owner's transparent mark: OUT
        m["_tree"] = a["owner_tree"]
        m["_owner_moved"] = (a["owner_tree"] != a["orig_tree"]
                             or a["owner_tier"] != a["orig_tier"])
        to = adj_by.get((m["_tree"], m["label"]))
        m["_adjusted"] = to is not None and to != a["owner_tier"]
        m["_tier"] = to or a["owner_tier"]
        surv.append(m)
    by_tab = {}
    for m in surv:
        m["_cat"] = categorize(m)
        by_tab.setdefault(m["_tree"], []).append(m)
    counts = Counter(m["_cat"] for m in surv)

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]
    legend_slide(prs, counts)

    listed = {t for t, _, _ in TREES}
    roster = list(TREES) + [(t, "888888", "(tab not in the roster — check the model)")
                            for t in sorted(by_tab) if t not in listed]

    total_placed = 0
    for tab, accent, desc in roster:
        rows = by_tab.get(tab, [])
        if not rows:
            continue
        slide = prs.slides.add_slide(blank)

        # Header strip.
        add_box(slide, 0, 0, SLIDE_W, HEADER_H, "", "2b2418", "2b2418",
                10, shape=MSO_SHAPE.RECTANGLE)
        title = slide.shapes.add_textbox(MARGIN, Emu(0), SLIDE_W - 2 * MARGIN, HEADER_H)
        tf = title.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = tab + "  "
        r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = rgb(accent)
        r2 = p.add_run(); r2.text = "%s   (%d techs)" % (desc, len(rows))
        r2.font.size = Pt(11); r2.font.color.rgb = rgb("d8cfc0")
        no_autofit(tf)

        top = HEADER_H + Emu(int(Inches(0.06)))
        area_top = top + TIERBAR_H + VGAP
        area_h = SLIDE_H - area_top - MARGIN
        rows_per_col = int((area_h + VGAP) // (BOX_H + VGAP))

        for ti, tier in enumerate(TIERS):
            cx = MARGIN + COL_W * ti
            add_box(slide, cx, top, COL_W - HGAP, TIERBAR_H, TIER_LABEL[tier],
                    TIER_FILL[tier], accent, 9, bold=True)
            members = sorted((m for m in rows if tier_of(m) == tier),
                             key=lambda m: (cost_of(m), m["label"]))
            if len(members) > SUB_COLS * rows_per_col:
                raise SystemExit("tier %s of %s holds %d boxes; capacity %d — "
                                 "shrink BOX_H" % (tier, tab, len(members),
                                                   SUB_COLS * rows_per_col))
            for i, m in enumerate(members):
                sc, sr = divmod(i, rows_per_col)
                x = cx + sc * (BOX_W + HGAP)
                y = area_top + sr * (BOX_H + VGAP)
                b = add_box(slide, x, y, BOX_W, BOX_H, m["label"],
                            CAT_FILL[m["_cat"]], accent, font_for(m["label"]))
                if m.get("_owner_moved"):
                    b.line.width = Pt(2.25)
                elif m.get("_adjusted"):
                    b.line.width = Pt(1.5)
                    b.line.dash_style = MSO_LINE_DASH.DASH
                total_placed += 1

    prs.save(OUT)
    print("wrote %s: %d slides (1 legend + trees), %d boxes / %d surviving rows"
          % (OUT, len(prs.slides._sldIdLst), total_placed, len(surv)))
    print("categories:", ", ".join("%s %d" % (k, v) for k, v in counts.most_common()))
    if total_placed != len(surv):
        print("MISMATCH: %d rows did not land on any slide" % (len(surv) - total_placed))
        sys.exit(1)


if __name__ == "__main__":
    build()
